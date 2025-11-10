import nbformat
from nbformat import NotebookNode
from PyPDF2 import PdfReader
import docx
import pptx
from typing import Union, Optional
import logging
from typing import List



logger = logging.getLogger(__name__)

class FileProcessor:
    def extract_code_from_ipynb(self, file_content: bytes) -> str:
        """Memory-efficient notebook processing"""
        try:
            nb = nbformat.reads(file_content.decode('utf-8'), as_version=4)
            output = []
            
            for idx, cell in enumerate(nb.cells):
                if cell.cell_type == 'code':
                    lines = cell.source.count('\n') + 1
                    output.append(f"\n=== Code Cell (Lines: {lines}) ===\n{cell.source}\n")
                    if len(output) > 10:
                        output.append("\n...[truncated due to size]")
                        break

            
            return "".join(output)
        except Exception as e:
            logger.error(f"Notebook processing failed: {str(e)}")
            return "Failed to process notebook"


    def extract_text_from_pdf(self, file_stream) -> str:
        """Extract text from PDF files"""
        try:
            reader = PdfReader(file_stream)
            return '\n'.join(page.extract_text() or '' for page in reader.pages)
        except Exception as e:
            logger.error(f"PDF extraction failed: {str(e)}")
            raise ValueError("Invalid PDF file")

    def extract_text_from_docx(self, file_stream) -> str:
        """Extract text from Word documents"""
        try:
            doc = docx.Document(file_stream)
            return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
            raise ValueError("Invalid Word document")

    def extract_text_from_pptx(self, file_stream) -> str:
        """Extract text from PowerPoint files"""
        try:
            prs = pptx.Presentation(file_stream)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            raise ValueError("Invalid PowerPoint file")

    def process_txt(self, file_content: bytes) -> str:
        """Process plain text files"""
        try:
            return file_content.decode('utf-8')
        except Exception as e:
            logger.error(f"Text processing failed: {str(e)}")
            raise ValueError("Invalid text file")